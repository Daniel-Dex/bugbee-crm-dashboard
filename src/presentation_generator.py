from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PINK = RGBColor(232, 84, 106)
DARK = RGBColor(28, 28, 46)
GRAY = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
PAPER = RGBColor(255, 248, 244)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(239, 68, 68)


def money(v: float) -> str:
    return f"R$ {v:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")


def pct(v: float | None) -> str:
    if v is None:
        return "n/d"
    return f"{v * 100:.1f}%".replace(".", ",")


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_kpi(slide, label, value, note, x, y):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.6), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(240, 217, 221)
    add_text(slide, label.upper(), x + 0.12, y + 0.12, 2.3, 0.18, 7.5, GRAY, True)
    add_text(slide, value, x + 0.12, y + 0.34, 2.3, 0.28, 17, DARK, True)
    add_text(slide, note, x + 0.12, y + 0.70, 2.3, 0.18, 7.5, GRAY)


def add_header(slide, kicker, title):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_text(slide, "BUGBEE", 0.55, 0.38, 2.0, 0.2, 10, PINK, True)
    add_text(slide, kicker.upper(), 0.55, 0.70, 3.2, 0.2, 8, GRAY, True)
    add_text(slide, title, 0.55, 0.98, 11.7, 0.55, 22, DARK, True)


def add_table(slide, rows, x, y, w, row_h=0.34):
    headers = ["Indicador", "Valor", "Nota"]
    widths = [w * 0.42, w * 0.23, w * 0.35]
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(row_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.color.rgb = DARK
    cursor = x
    for header, width in zip(headers, widths):
        add_text(slide, header.upper(), cursor + 0.08, y + 0.08, width - 0.12, 0.14, 7, WHITE, True)
        cursor += width
    for i, row in enumerate(rows):
        yy = y + row_h * (i + 1)
        shape = slide.shapes.add_shape(1, Inches(x), Inches(yy), Inches(w), Inches(row_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(253, 240, 243)
        shape.line.color.rgb = RGBColor(240, 217, 221)
        cursor = x
        for value, width in zip(row, widths):
            add_text(slide, str(value), cursor + 0.08, yy + 0.08, width - 0.12, 0.14, 7.5, DARK, i == 0)
            cursor += width


def generate_presentation(output_path: Path, period_label: str, metrics: dict, insights: list[str]) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK
    add_text(slide, "BUGBEE", 0.65, 0.72, 3.5, 0.4, 20, WHITE, True)
    add_text(slide, "Performance semanal de e-commerce", 0.65, 1.70, 8.5, 0.7, 28, WHITE, True)
    add_text(slide, period_label, 0.65, 2.42, 4.0, 0.3, 13, PINK, True)
    add_kpi(slide, "Faturamento", money(metrics.get("faturamento", 0)), "semana analisada", 0.65, 5.75)
    add_kpi(slide, "Pedidos", f"{metrics.get('pedidos', 0):.0f}", "pedidos capturados", 3.45, 5.75)
    add_kpi(slide, "Ticket médio", money(metrics.get("ticket_medio", 0)), "faturamento / pedidos", 6.25, 5.75)
    add_kpi(slide, "Itens", f"{metrics.get('itens_vendidos', 0):.0f}", "peças vendidas", 9.05, 5.75)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Resumo executivo", "Principais leituras da semana")
    for i, item in enumerate(insights[:5]):
        add_text(slide, f"{i + 1}. {item}", 0.75, 1.85 + i * 0.62, 11.5, 0.35, 14, DARK)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Indicadores", "Visão geral de vendas")
    add_kpi(slide, "Faturamento", money(metrics.get("faturamento", 0)), "R$ total capturado", 0.65, 2.05)
    add_kpi(slide, "Pedidos", f"{metrics.get('pedidos', 0):.0f}", "quantidade", 3.45, 2.05)
    add_kpi(slide, "Ticket médio", money(metrics.get("ticket_medio", 0)), "R$ / pedido", 6.25, 2.05)
    add_kpi(slide, "Peças/pedido", f"{metrics.get('pecas_por_pedido', 0):.2f}".replace(".", ","), "itens / pedido", 9.05, 2.05)
    add_text(slide, f"Desconto médio: {pct(metrics.get('desconto_medio'))}", 0.70, 3.55, 4.5, 0.3, 15, DARK, True)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "DRE simplificado", "A leitura financeira fica em tabela para reconciliar com a planilha base")
    rows = [
        ("Faturamento", money(metrics.get("faturamento", 0)), "pedidos API"),
        ("Pedidos", f"{metrics.get('pedidos', 0):.0f}", "cabeçalho pedido"),
        ("Peças vendidas", f"{metrics.get('itens_vendidos', 0):.0f}", "itens de nota fiscal"),
        ("Ticket médio", money(metrics.get("ticket_medio", 0)), "faturamento / pedidos"),
        ("Peças por pedido", f"{metrics.get('pecas_por_pedido', 0):.2f}".replace(".", ","), "peças / pedidos"),
        ("Desconto médio", pct(metrics.get("desconto_medio")), "quando valor de desconto está disponível"),
    ]
    add_table(slide, rows, 0.75, 1.85, 11.5)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Próximos passos", "O que validar antes de colocar a automação no piloto semanal")
    next_steps = [
        "Reconciliar pedidos API vs pedidos faturados da planilha Base Análise Magazord.",
        "Validar se notas fiscais são a melhor fonte para peças, produtos e descontos.",
        "Liberar ou substituir endpoints 403 de produtos completos e canais de venda.",
        "Conectar mídia/funil por fonte complementar, caso não estejam na Magazord.",
    ]
    for i, item in enumerate(next_steps):
        add_text(slide, f"{i + 1}. {item}", 0.75, 1.85 + i * 0.62, 11.5, 0.35, 14, DARK)

    prs.save(output_path)
